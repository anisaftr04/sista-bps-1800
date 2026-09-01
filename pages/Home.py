import streamlit as st
import pandas as pd
from supabase import create_client, Client
import plotly.express as px
from datetime import datetime
import os
import base64
from docx import Document
from pptx import Presentation
import tempfile


# ============================================================
# KONEKSI SUPABASE
# ============================================================

url = st.secrets["SUPABASE_URL"]
key = st.secrets["SUPABASE_KEY"]

supabase: Client = create_client(url, key)


# ============================================================
# FUNGSI JENIS FILE
# ============================================================

def get_jenis_file(nama_file):

    ext = os.path.splitext(nama_file)[1].lower()

    if ext == ".pdf":
        return "PDF"

    elif ext in [".xlsx", ".xls"]:
        return "Excel"

    elif ext == ".docx":
        return "Word"

    elif ext == ".pptx":
        return "PPT"

    else:
        return "Lainnya"


# ============================================================
# MEMBACA DATA DOKUMEN
# ============================================================

data = []

response = (
    supabase
    .table("dokumen")
    .select("*")
    .order("created_at", desc=True)
    .execute()
)


for doc in response.data:

    data.append(
        {
            "Nama": doc["nama_file"],
            "Kategori": doc["kategori"],
            "Jenis": get_jenis_file(doc["nama_file"]),
            "Ukuran": "-",
            "Tanggal": doc["created_at"],
            "Path": doc["path_file"]
        }
    )


df = pd.DataFrame(
    data,
    columns=[
        "Nama",
        "Kategori",
        "Jenis",
        "Ukuran",
        "Tanggal",
        "Path"
    ]
)


if df.empty:
    st.warning("Belum ada dokumen di database.")


# ============================================================
# HEADER
# ============================================================

jam = datetime.now().strftime("%H:%M")


# ------------------------------------------------------------
# LOGO BPS
# ------------------------------------------------------------

with open("logo_bps.png", "rb") as file:

    logo_base64 = base64.b64encode(
        file.read()
    ).decode()


st.html(
    f"""
    <div class="bps-brand">

        <img
            src="data:image/png;base64,{logo_base64}"
            class="bps-logo"
        >

        <div class="bps-text">

            <div class="bps-title">
                Badan Pusat Statistik
            </div>

            <div class="bps-province">
                Provinsi Lampung
            </div>

        </div>

    </div>
    """
)


# ------------------------------------------------------------
# HEADER SISTA
# ------------------------------------------------------------

st.html(
    f"""
    <div class="header">

        <h1>SISTA</h1>

        <h3>
            Sistem Informasi Statistik Sosial
        </h3>

        <p>
            BPS Provinsi Lampung
        </p>

        <p>
            🕒 {jam} WIB
        </p>

    </div>
    """
)


st.write("")


# ============================================================
# KPI
# ============================================================

col1, col2, col3, col4 = st.columns(4)


with col1:

    st.html(
        f"""
        <div class="card biru">

            <h2>📄</h2>

            <h1>
                {len(df)}
            </h1>

            <p>
                Total Dokumen
            </p>

        </div>
        """
    )


with col2:

    st.html(
        f"""
        <div class="card merah">

            <h2>📕</h2>

            <h1>
                {(df["Jenis"] == "PDF").sum()}
            </h1>

            <p>
                File PDF
            </p>

        </div>
        """
    )


with col3:

    st.html(
        f"""
        <div class="card hijau">

            <h2>📗</h2>

            <h1>
                {(df["Jenis"] == "Excel").sum()}
            </h1>

            <p>
                File Excel
            </p>

        </div>
        """
    )


with col4:

    st.html(
        f"""
        <div class="card orange">

            <h2>📊</h2>

            <h1>
                {(df["Jenis"] == "PPT").sum()}
            </h1>

            <p>
                File PPT
            </p>

        </div>
        """
    )


st.write("")


# ============================================================
# GRAFIK
# ============================================================

k1, k2 = st.columns([2, 1])


with k1:

    if not df.empty:

        fig = px.bar(
            df,
            x="Kategori",
            color="Kategori",
            title="📊 Jumlah Dokumen per Kategori",
            color_discrete_sequence=[
                "#005BAC",
                "#F7941D",
                "#4CAF50",
                "#8E44AD",
                "#E74C3C"
            ]
        )

        fig.update_layout(
            plot_bgcolor="white",
            paper_bgcolor="white",
            title_font_size=22,
            title_x=0.02,
            xaxis_title="Kategori",
            yaxis_title="Jumlah Dokumen",
            showlegend=False
        )

        st.plotly_chart(
            fig,
            use_container_width=True
        )


with k2:

    if not df.empty:

        pie = px.pie(
            df,
            names="Jenis",
            title="📁 Komposisi Jenis Dokumen",
            color_discrete_sequence=[
                "#005BAC",
                "#F7941D",
                "#4CAF50",
                "#8E44AD"
            ]
        )

        pie.update_layout(
            paper_bgcolor="white",
            plot_bgcolor="white",
            title_font_size=20,
            title_x=0.05
        )

        pie.update_traces(
            textinfo="percent+label",
            hole=0.45
        )

        st.plotly_chart(
            pie,
            use_container_width=True
        )


st.write("")


# ============================================================
# TABEL DOKUMEN
# ============================================================

st.subheader("📂 Dokumen")

st.caption(
    "Dokumen dikelompokkan berdasarkan jenis file"
)


jenis_file = [
    "PDF",
    "Excel",
    "Word",
    "PPT"
]


for jenis in jenis_file:

    data_jenis = df[
        df["Jenis"] == jenis
    ]

    with st.expander(
        f"📁 {jenis} ({len(data_jenis)})",
        expanded=False
    ):

        if data_jenis.empty:

            st.info(
                f"Tidak ada file {jenis}."
            )

            continue


        for _, row in data_jenis.iterrows():

            st.html(
                f"""
                <div class="dokumen-card">

                    <h4>
                        📄 {row['Nama']}
                    </h4>

                    <p>
                        📂 <b>Jenis</b> : {row['Jenis']}
                    </p>

                    <p>
                        📅 <b>Tanggal</b> : {row['Tanggal']}
                    </p>

                    <p>
                        📦 <b>Ukuran</b> : {row['Ukuran']}
                    </p>

                </div>
                """
            )


            col1, col2, col3 = st.columns(
                [6, 1, 1]
            )


            # ------------------------------------------------
            # LIHAT
            # ------------------------------------------------

            with col2:

                if st.button(
                    "👁 Lihat",
                    key=f"lihat_{row['Nama']}"
                ):

                    st.session_state[
                        "file_dipilih"
                    ] = row["Path"]


            # ------------------------------------------------
            # DOWNLOAD
            # ------------------------------------------------

            with col3:

                try:

                    file_data = (
                        supabase
                        .storage
                        .from_("dokumen")
                        .download(row["Path"])
                    )

                    st.download_button(
                        label="⬇ Download",
                        data=file_data,
                        file_name=row["Nama"],
                        key=f"download_{row['Nama']}"
                    )

                except Exception as e:

                    st.error(
                        f"Gagal download: {e}"
                    )


            # ------------------------------------------------
            # HAPUS
            # ------------------------------------------------

            if st.session_state.get(
                "is_admin",
                False
            ):

                col_hapus, _ = st.columns(
                    [1, 5]
                )

                with col_hapus:

                    if st.button(
                        "🗑 Hapus",
                        key=f"hapus_{row['Nama']}"
                    ):

                        st.session_state[
                            "hapus_file"
                        ] = row["Path"]


            st.divider()


# ============================================================
# KONFIRMASI HAPUS
# ============================================================

if "hapus_file" in st.session_state:

    path_file = st.session_state[
        "hapus_file"
    ]

    st.warning(
        "Yakin ingin menghapus dokumen ini?"
    )

    col1, col2 = st.columns(2)


    with col1:

        if st.button(
            "✅ Ya, Hapus",
            key="confirm_hapus"
        ):

            try:

                # Hapus file dari Storage
                supabase.storage \
                    .from_("dokumen") \
                    .remove([path_file])


                # Hapus data dari tabel
                (
                    supabase
                    .table("dokumen")
                    .delete()
                    .eq("path_file", path_file)
                    .execute()
                )


                del st.session_state[
                    "hapus_file"
                ]


                st.success(
                    "Dokumen berhasil dihapus."
                )

                st.rerun()


            except Exception as e:

                st.error(
                    f"Gagal menghapus dokumen: {e}"
                )


    with col2:

        if st.button(
            "❌ Batal",
            key="batal_hapus"
        ):

            del st.session_state[
                "hapus_file"
            ]

            st.rerun()


# ============================================================
# PREVIEW DOKUMEN
# ============================================================

if "file_dipilih" in st.session_state:

    st.markdown("---")

    st.subheader(
        "👁 Preview Dokumen"
    )


    try:

        storage_path = st.session_state[
            "file_dipilih"
        ]


        # ----------------------------------------------------
        # DOWNLOAD FILE
        # ----------------------------------------------------

        file_bytes = (
            supabase
            .storage
            .from_("dokumen")
            .download(storage_path)
        )


        suffix = os.path.splitext(
            storage_path
        )[1]


        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=suffix
        ) as tmp:

            tmp.write(file_bytes)

            local_path = tmp.name


        # ====================================================
        # PDF
        # ====================================================

        if suffix.lower() == ".pdf":

            with open(
                local_path,
                "rb"
            ) as pdf_file:

                base64_pdf = base64.b64encode(
                    pdf_file.read()
                ).decode("utf-8")


            st.html(
                f"""
                <iframe
                    src="data:application/pdf;base64,{base64_pdf}"
                    width="100%"
                    height="700px"
                    style="border:none;"
                >
                </iframe>
                """
            )


        # ====================================================
        # EXCEL
        # ====================================================

        elif suffix.lower() in [
            ".xlsx",
            ".xls"
        ]:

            excel = pd.read_excel(
                local_path
            )

            st.dataframe(
                excel,
                use_container_width=True,
                height=600
            )


        # ====================================================
        # WORD
        # ====================================================

        elif suffix.lower() == ".docx":

            doc = Document(
                local_path
            )

            isi = ""

            for paragraf in doc.paragraphs:

                isi += (
                    paragraf.text
                    + "\n"
                )


            st.text_area(
                "Isi Dokumen",
                isi,
                height=600
            )


        # ====================================================
        # POWERPOINT
        # ====================================================

        elif suffix.lower() == ".pptx":

            prs = Presentation(
                local_path
            )


            for i, slide in enumerate(
                prs.slides
            ):

                with st.expander(
                    f"Slide {i+1}",
                    expanded=True
                ):

                    ada_teks = False


                    for shape in slide.shapes:

                        if hasattr(
                            shape,
                            "text"
                        ):

                            teks = (
                                shape.text
                                .strip()
                            )


                            if teks:

                                ada_teks = True

                                st.write(
                                    teks
                                )


                    if not ada_teks:

                        st.info(
                            "Slide ini tidak memiliki teks."
                        )


        else:

            st.info(
                "Preview untuk jenis file ini belum didukung."
            )


    except Exception as e:

        st.error(
            f"Gagal membuka dokumen: {e}"
        )